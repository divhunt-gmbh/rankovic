#!/usr/bin/env python3
"""
TALIJA by Ranković - PowerPoint Presentation Generator
Generates a bilingual (Serbian + Chinese) presentation for Chinese market
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Colors
GOLD = RgbColor(201, 162, 39)
DARK = RgbColor(13, 13, 13)
DARK_SOFT = RgbColor(30, 30, 30)
BEIGE = RgbColor(245, 240, 230)
WHITE = RgbColor(255, 255, 255)

# Base path for images
BASE_PATH = os.path.dirname(os.path.abspath(__file__))
IMAGES_PATH = os.path.join(BASE_PATH, "images")

def add_dark_background(slide):
    """Add dark background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_light_background(slide):
    """Add light/beige background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BEIGE
    background.line.fill.background()
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_gold_background(slide):
    """Add gold background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = GOLD
    background.line.fill.background()
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box to slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image to slide if it exists"""
    full_path = os.path.join(BASE_PATH, image_path)
    if os.path.exists(full_path):
        if width and height:
            return slide.shapes.add_picture(full_path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            return slide.shapes.add_picture(full_path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            return slide.shapes.add_picture(full_path, Inches(left), Inches(top), height=Inches(height))
        else:
            return slide.shapes.add_picture(full_path, Inches(left), Inches(top))
    else:
        # Add placeholder rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width or 4), Inches(height or 3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_SOFT
        shape.line.color.rgb = GOLD
        return shape

def add_slide_number(slide, number, total=14, color=WHITE):
    """Add slide number to bottom right"""
    add_text_box(slide, 11.5, 6.9, 1.5, 0.4, f"{number:02d} / {total}",
                 font_size=10, color=color, align=PP_ALIGN.RIGHT)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ==================== SLIDE 1: Title ====================
    slide1 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide1)

    # Brand name
    add_text_box(slide1, 0, 2, 13.333, 1, "TALIJA",
                 font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide1, 0, 2.9, 13.333, 0.5, "by Ranković",
                 font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # Gold line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(3.5), Inches(1.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Tagline
    add_text_box(slide1, 0, 3.8, 13.333, 0.5, "Porodično Nasleđe",
                 font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide1, 0, 4.3, 13.333, 0.5, "家族传承与当代表达的融合",
                 font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide1, 0, 5, 13.333, 0.5, "Premium Srpska Rakija · 塞尔维亚优质拉基亚",
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide1, 0, 5.5, 13.333, 0.4, "Porodično nasleđe pretočeno u savremeni izraz",
                 font_size=14, color=RgbColor(200, 200, 200), align=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide1, 0, 6.3, 13.333, 0.4, "Destilerija Ranković · Est. 2022 · Lazarevac, Srbija",
                 font_size=12, color=RgbColor(150, 150, 150), align=PP_ALIGN.CENTER)

    add_slide_number(slide1, 1)

    # ==================== SLIDE 2: Serbia & Rakija ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_light_background(slide2)

    # Title
    add_text_box(slide2, 0, 0.5, 13.333, 0.7, "Srbija – Zemlja Rakije",
                 font_size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide2, 0, 1.1, 13.333, 0.5, "塞尔维亚 - 拉基亚之乡",
                 font_size=24, color=DARK, align=PP_ALIGN.CENTER)

    # Image
    add_image_safe(slide2, "images/viber_slika_2025-12-08_16-15-36-688.jpg", 0.8, 2, 5.5, 4)

    # Features list
    features = [
        ("🍇 Decenijska tradicija · 数十年的传统", "Porodična proizvodnja rakije kroz generacije"),
        ("🏠 Porodična tradicija · 家族传统", "Svaka porodica ima svoju recepturu"),
        ("🤝 Simbol gostoprimstva · 好客的象征", "Rakija se služi gostima kao znak dobrodošlice"),
        ("🌍 Geografski zaštićen proizvod · 地理标志保护产品", "Autentičan evropski proizvod"),
    ]

    y_pos = 2.2
    for title, desc in features:
        add_text_box(slide2, 6.8, y_pos, 5.5, 0.4, title,
                     font_size=14, bold=True, color=DARK)
        add_text_box(slide2, 6.8, y_pos + 0.35, 5.5, 0.4, desc,
                     font_size=11, color=RgbColor(80, 80, 80))
        y_pos += 1.0

    add_slide_number(slide2, 2, color=DARK)

    # ==================== SLIDE 3: Destilerija Ranković ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide3)

    # Title
    add_text_box(slide3, 0, 0.5, 13.333, 0.7, "Destilerija Ranković",
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide3, 0, 1.1, 13.333, 0.5, "兰科维奇酿酒厂",
                 font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # Image
    add_image_safe(slide3, "images/new/WhatsApp Image 2025-12-16 at 9.35.29 PM (1).jpeg", 0.8, 2, 5.5, 4)

    # Content
    add_text_box(slide3, 6.8, 2, 5.5, 0.6, "Znanje Koje Se Ne Prekida",
                 font_size=20, bold=True, color=GOLD)
    add_text_box(slide3, 6.8, 2.4, 5.5, 0.4, "不曾中断的技艺传承",
                 font_size=14, color=WHITE)

    add_text_box(slide3, 6.8, 3, 5.5, 1.4,
                 "Znanje o pečenju rakije u porodici Ranković prenosi se kroz tri generacije. Prvi je ovaj zanat započeo deda. Danas rakiju proizvode otac i sin zajedno.",
                 font_size=12, color=WHITE)
    add_text_box(slide3, 6.8, 4.3, 5.5, 0.8,
                 "兰科维奇家族的蒸馏技艺已传承三代，如今由父子共同酿造。",
                 font_size=11, color=RgbColor(180, 180, 180))

    # Stats
    add_text_box(slide3, 6.8, 5.3, 1.7, 0.6, "3", font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide3, 6.8, 5.8, 1.7, 0.3, "Generacije · 代", font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide3, 8.6, 5.3, 1.7, 0.6, "10", font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide3, 8.6, 5.8, 1.7, 0.3, "Zlatnih medalja · 金奖", font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide3, 10.4, 5.3, 1.7, 0.6, "4", font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide3, 10.4, 5.8, 1.7, 0.3, "Vrste rakije · 品种", font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    add_slide_number(slide3, 3)

    # ==================== SLIDE 4: Philosophy ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_gold_background(slide4)

    # Title
    add_text_box(slide4, 0, 0.8, 13.333, 0.7, "Tradicija Vođena Znanjem",
                 font_size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide4, 0, 1.4, 13.333, 0.5, "以知识引导的传统",
                 font_size=24, color=DARK, align=PP_ALIGN.CENTER)

    # Quote
    add_text_box(slide4, 1.5, 2.3, 10.333, 1.2,
                 '"U porodici Ranković znanje o destilaciji ne smatra se ličnom veštinom, već obavezom prema precima i odgovornošću prema generacijama koje dolaze."',
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide4, 1.5, 3.5, 10.333, 0.8,
                 "在兰科维奇家族中，蒸馏技艺并非个人能力的体现，而是一种对祖辈的责任，以及对未来世代的承诺。",
                 font_size=15, color=DARK, align=PP_ALIGN.CENTER)

    # Gold line (dark on gold bg)
    line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(4.5), Inches(1.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK
    line.line.fill.background()

    add_text_box(slide4, 1.5, 4.9, 10.333, 0.8,
                 "Pravi kvalitet se stalno potvrđuje učenjem i usavršavanjem. Trajna vrednost gradi se postepeno, kroz dosledan rad.",
                 font_size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide4, 1.5, 5.7, 10.333, 0.5,
                 "真正的品质需要通过持续学习与精进不断验证。真正的价值来自循序渐进的坚持。",
                 font_size=12, color=DARK, align=PP_ALIGN.CENTER)

    add_slide_number(slide4, 4, color=DARK)

    # ==================== SLIDE 5: Four Pillars ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide5)

    # Title
    add_text_box(slide5, 0, 0.5, 13.333, 0.7, "Četiri Stuba Kvaliteta",
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide5, 0, 1.1, 13.333, 0.5, "质量四大支柱",
                 font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # Four pillars
    pillars = [
        ("🍎", "Čisto Voće", "纯净水果", "100% prirodno voće bez aditiva", "100%天然水果，无添加剂"),
        ("🔥", "Dvostruka Destilacija", "双重蒸馏", "Tradicionalne metode", "传统工艺"),
        ("❤️", "Sa Ljubavlju", "用心酿造", "Ručna proizvodnja, mala serija", "手工制作，小批量生产"),
        ("🏆", "Premium Kvalitet", "优质品质", "Bez kompromisa", "绝不妥协"),
    ]

    x_pos = 0.8
    for emoji, title_sr, title_cn, desc_sr, desc_cn in pillars:
        # Card background
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2), Inches(2.8), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_SOFT
        card.line.color.rgb = GOLD

        add_text_box(slide5, x_pos, 2.3, 2.8, 0.6, emoji, font_size=40, align=PP_ALIGN.CENTER)
        add_text_box(slide5, x_pos, 3, 2.8, 0.4, title_sr, font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide5, x_pos, 3.4, 2.8, 0.4, title_cn, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide5, x_pos + 0.1, 4, 2.6, 0.6, desc_sr, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide5, x_pos + 0.1, 4.5, 2.6, 0.5, desc_cn, font_size=10, color=RgbColor(180, 180, 180), align=PP_ALIGN.CENTER)

        x_pos += 3.1

    add_slide_number(slide5, 5)

    # ==================== SLIDE 6: Collection Overview ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_light_background(slide6)

    # Title
    add_text_box(slide6, 0, 0.5, 13.333, 0.7, "TALIJA Kolekcija",
                 font_size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide6, 0, 1.1, 13.333, 0.5, "塔利亚系列",
                 font_size=24, color=DARK, align=PP_ALIGN.CENTER)

    add_text_box(slide6, 0, 1.8, 13.333, 0.5, "Četiri ukusa, jedna priča · 四种口味，一个故事",
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)

    # Four products
    products = [
        ("🟣", "Šljiva", "李子"),
        ("🍏", "Jabuka", "苹果"),
        ("🍐", "Kruška", "梨子"),
        ("🟡", "Dunja", "榅桲"),
    ]

    x_pos = 0.8
    for emoji, name_sr, name_cn in products:
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2.8), Inches(2.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RgbColor(220, 220, 220)

        add_text_box(slide6, x_pos, 3.1, 2.8, 0.8, emoji, font_size=48, align=PP_ALIGN.CENTER)
        add_text_box(slide6, x_pos, 4, 2.8, 0.4, name_sr, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide6, x_pos, 4.4, 2.8, 0.3, name_cn, font_size=12, color=DARK, align=PP_ALIGN.CENTER)

        x_pos += 3.1

    add_text_box(slide6, 0, 5.8, 13.333, 0.4, "TALIJA je naša ljubavna pesma srpskom voću.",
                 font_size=14, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide6, 0, 6.2, 13.333, 0.4, "塔利亚是我们对塞尔维亚水果的爱的颂歌。",
                 font_size=12, color=RgbColor(100, 100, 100), align=PP_ALIGN.CENTER)

    add_slide_number(slide6, 6, color=DARK)

    # ==================== SLIDE 7: Šljiva ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide7)

    # Image
    add_image_safe(slide7, "images/viber_slika_2025-12-08_16-15-33-501.jpg", 0.5, 0.8, 5.5, 5.8)

    # Content
    add_text_box(slide7, 6.5, 1, 6, 0.3, "Srce Destilerije · 酒坊的核心之作",
                 font_size=10, color=GOLD)
    add_text_box(slide7, 6.5, 1.5, 6, 0.6, "TALIJA Šljiva",
                 font_size=32, bold=True, color=WHITE)
    add_text_box(slide7, 6.5, 2.1, 6, 0.4, "塔利亚李子白兰地",
                 font_size=18, color=WHITE)

    add_text_box(slide7, 6.5, 2.7, 6, 0.5, "🏆 Zlatna medalja · Novosadski sajam 2025",
                 font_size=14, color=GOLD)

    add_text_box(slide7, 6.5, 3.4, 6, 1.4,
                 "Centralni proizvod destilerije i najviši izraz znanja porodice Ranković. Talija nije rakija od jedne sorte – ona je pažljivo razvijen blend više destilata.",
                 font_size=12, color=WHITE)
    add_text_box(slide7, 6.5, 4.8, 6, 1,
                 "酒坊的核心产品，代表了兰科维奇家族技艺与经验的最高水平。TALI娅并非单一品种白兰地，而是一款精心调配而成的复合酒。",
                 font_size=11, color=RgbColor(180, 180, 180))

    add_text_box(slide7, 6.5, 6, 6, 0.4, "Ravnoteža daje dubinu, stabilnost i vrednost · 平衡赋予深度与价值",
                 font_size=10, color=RgbColor(120, 120, 120))

    add_slide_number(slide7, 7)

    # ==================== SLIDE 8: Jabuka ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide8)

    # Content (left side)
    add_text_box(slide8, 0.8, 1, 6, 0.3, "TALIJA COLLECTION",
                 font_size=10, color=GOLD)
    add_text_box(slide8, 0.8, 1.5, 6, 0.6, "TALIJA Jabuka",
                 font_size=32, bold=True, color=WHITE)
    add_text_box(slide8, 0.8, 2.1, 6, 0.4, "塔利亚苹果白兰地",
                 font_size=18, color=WHITE)

    add_text_box(slide8, 0.8, 2.8, 6, 0.5, '"Jutarnja Svetlost" · "晨曦之光"',
                 font_size=18, color=GOLD)

    add_text_box(slide8, 0.8, 3.6, 5.5, 1.2,
                 "Sveža, živahna aroma zelenih i crvenih jabuka sa citruznim akcentima. Ukus je balansiran – slatko-kiselkasto, sa blagom začinskom notom.",
                 font_size=13, color=WHITE)
    add_text_box(slide8, 0.8, 4.8, 5.5, 1,
                 "新鲜活泼的青苹果和红苹果香气，带有柑橘的点缀。口感平衡——酸甜适中，带有淡淡的香料味。",
                 font_size=11, color=RgbColor(180, 180, 180))

    add_text_box(slide8, 0.8, 6, 5.5, 0.4, "Osvežavajuća i elegantna · 清爽优雅",
                 font_size=10, color=RgbColor(120, 120, 120))

    # Image (right side)
    add_image_safe(slide8, "images/viber_slika_2025-12-08_16-15-33-284.jpg", 7.3, 0.8, 5.5, 5.8)

    add_slide_number(slide8, 8)

    # ==================== SLIDE 9: Kruška ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide9)

    # Image
    add_image_safe(slide9, "images/viber_image_2025-12-16_21-44-44-133.jpg", 0.5, 0.8, 5.5, 5.8)

    # Content
    add_text_box(slide9, 6.5, 1, 6, 0.3, "TALIJA COLLECTION",
                 font_size=10, color=GOLD)
    add_text_box(slide9, 6.5, 1.5, 6, 0.6, "TALIJA Kruška",
                 font_size=32, bold=True, color=WHITE)
    add_text_box(slide9, 6.5, 2.1, 6, 0.4, "塔利亚梨子白兰地",
                 font_size=18, color=WHITE)

    add_text_box(slide9, 6.5, 2.8, 6, 0.5, '"Kristalna Elegancija" · "水晶般的优雅"',
                 font_size=18, color=GOLD)

    add_text_box(slide9, 6.5, 3.6, 6, 1.2,
                 "Mirisna, cvetna aroma odabranih sorti krušaka koja otvara čula. Ukus je svilenkast, mekan, sa fino izbalansiranom slatkoćom i diskretnom kiselošću.",
                 font_size=13, color=WHITE)
    add_text_box(slide9, 6.5, 4.8, 6, 1,
                 "威廉斯梨的芬芳花香，唤醒感官。口感如丝般柔滑，甜度平衡，带有微妙的酸度。",
                 font_size=11, color=RgbColor(180, 180, 180))

    add_text_box(slide9, 6.5, 6, 6, 0.4, "Pažljiv odabir sorti · 精选品种",
                 font_size=10, color=RgbColor(120, 120, 120))

    add_slide_number(slide9, 9)

    # ==================== SLIDE 10: Dunja ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide10)

    # Content (left side)
    add_text_box(slide10, 0.8, 1, 6, 0.3, "TALIJA COLLECTION",
                 font_size=10, color=GOLD)
    add_text_box(slide10, 0.8, 1.5, 6, 0.6, "TALIJA Dunja",
                 font_size=32, bold=True, color=WHITE)
    add_text_box(slide10, 0.8, 2.1, 6, 0.4, "塔利亚榅桲白兰地",
                 font_size=18, color=WHITE)

    add_text_box(slide10, 0.8, 2.8, 6, 0.5, '"Zlatna Pesma" · "金色之歌"',
                 font_size=18, color=GOLD)

    add_text_box(slide10, 0.8, 3.6, 5.5, 1.2,
                 "Bogata, složena aroma dunje sa cvetnim notama kamilice i toplim mednim tonovima. Završnica je duga, zlatna, aromatična.",
                 font_size=13, color=WHITE)
    add_text_box(slide10, 0.8, 4.8, 5.5, 1,
                 "榅桲的浓郁复杂香气，带有洋甘菊和烤杏的花香。余味悠长，金色，芳香四溢。",
                 font_size=11, color=RgbColor(180, 180, 180))

    add_text_box(slide10, 0.8, 6, 5.5, 0.4, "Retka i dragocena · 稀有珍贵",
                 font_size=10, color=RgbColor(120, 120, 120))

    # Image (right side)
    add_image_safe(slide10, "images/viber_slika_2025-12-08_16-15-33-077.jpg", 7.3, 0.8, 5.5, 5.8)

    add_slide_number(slide10, 10)

    # ==================== SLIDE 11: Tihi Luksuz ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_gold_background(slide11)

    # Title
    add_text_box(slide11, 0, 0.5, 13.333, 0.7, "Tihi Luksuz",
                 font_size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide11, 0, 1.1, 13.333, 0.5, "低调而内敛的奢华",
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)

    # Quote
    add_text_box(slide11, 1, 1.8, 11.333, 0.8,
                 "Talija svoju vrednost ne gradi kroz upadljivu promociju, već kroz poreklo, proces i priznanja.",
                 font_size=14, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide11, 1, 2.5, 11.333, 0.5,
                 "TALI娅的价值并不依赖张扬的宣传，而体现在其来源、工艺与获得的认可之中。",
                 font_size=11, color=RgbColor(60, 60, 60), align=PP_ALIGN.CENTER)

    # Reasons - left column
    reasons_left = [
        ("🏆 10 zlatnih medalja · 十枚金奖", "Novosadski sajam 2025 · 诺维萨德农博会"),
        ("🔒 Ograničena proizvodnja · 限量生产", "Potpuna kontrola kvaliteta · 全面品控"),
        ("🎁 Premium poklon · 高端礼品", "Gravirane čašice · 定制雕刻酒杯"),
    ]

    y_pos = 3.2
    for title, desc in reasons_left:
        add_text_box(slide11, 0.8, y_pos, 5.5, 0.4, title, font_size=13, bold=True, color=DARK)
        add_text_box(slide11, 0.8, y_pos + 0.35, 5.5, 0.4, desc, font_size=10, color=RgbColor(60, 60, 60))
        y_pos += 1.0

    # Reasons - right column
    reasons_right = [
        ("🌿 Prirodni proizvod · 天然产品", "100% voće, bez aditiva · 100%水果"),
        ("🤝 Dugoročna partnerstva · 长期合作", "Stabilnost i poverenje · 稳定与信任"),
        ("🌍 Autentičan proizvod · 正宗产品", "Iz srca Srbije · 来自塞尔维亚"),
    ]

    y_pos = 3.2
    for title, desc in reasons_right:
        add_text_box(slide11, 7, y_pos, 5.5, 0.4, title, font_size=13, bold=True, color=DARK)
        add_text_box(slide11, 7, y_pos + 0.35, 5.5, 0.4, desc, font_size=10, color=RgbColor(60, 60, 60))
        y_pos += 1.0

    add_slide_number(slide11, 11, color=DARK)

    # ==================== SLIDE 12: Cooperation ====================
    slide12 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide12)

    # Title
    add_text_box(slide12, 0, 0.5, 13.333, 0.7, "Mogućnosti Saradnje",
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide12, 0, 1.1, 13.333, 0.5, "合作机会",
                 font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    # Cooperation options
    options = [
        ("🤝", "Ekskluzivna Distribucija", "独家经销", "Ekskluzivna prava za regione", "区域独家经销权"),
        ("🏪", "Uvoz i Veleprodaja", "进口批发", "Direktan uvoz iz Srbije", "从塞尔维亚直接进口"),
        ("🍽️", "HoReCa", "酒店餐饮", "Hoteli, restorani, barovi", "酒店、餐厅、酒吧"),
        ("🎁", "Poklon Tržište", "礼品市场", "Premium pokloni i setovi", "高端礼品和套装"),
    ]

    x_pos = 0.8
    for emoji, title_sr, title_cn, desc_sr, desc_cn in options:
        card = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2), Inches(2.8), Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_SOFT
        card.line.color.rgb = GOLD

        add_text_box(slide12, x_pos, 2.3, 2.8, 0.6, emoji, font_size=36, align=PP_ALIGN.CENTER)
        add_text_box(slide12, x_pos, 3, 2.8, 0.4, title_sr, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide12, x_pos, 3.4, 2.8, 0.4, title_cn, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide12, x_pos + 0.1, 4, 2.6, 0.5, desc_sr, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide12, x_pos + 0.1, 4.4, 2.6, 0.5, desc_cn, font_size=9, color=RgbColor(180, 180, 180), align=PP_ALIGN.CENTER)

        x_pos += 3.1

    add_slide_number(slide12, 12)

    # ==================== SLIDE 13: Contact ====================
    slide13 = prs.slides.add_slide(blank_layout)
    add_light_background(slide13)

    # Title
    add_text_box(slide13, 0, 0.8, 13.333, 0.7, "Kontakt",
                 font_size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide13, 0, 1.4, 13.333, 0.5, "联系方式",
                 font_size=24, color=DARK, align=PP_ALIGN.CENTER)

    # Contact info
    contacts = [
        ("📍", "Adresa · 地址", "Medoševački Put 2a\nLazarevac, Srbija"),
        ("📞", "Telefon · 电话", "+381 65 383 00 10"),
        ("✉️", "Email · 邮箱", "destilerijarankovic@gmail.com"),
    ]

    x_pos = 1.5
    for emoji, title, info in contacts:
        add_text_box(slide13, x_pos, 2.8, 3.5, 0.6, emoji, font_size=36, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide13, x_pos, 3.5, 3.5, 0.4, title, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide13, x_pos, 4, 3.5, 1, info, font_size=13, color=DARK, align=PP_ALIGN.CENTER)
        x_pos += 3.8

    # Website
    add_text_box(slide13, 0, 5.5, 13.333, 0.5, "🌐 rakijatalija.rs",
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)

    add_slide_number(slide13, 13, color=DARK)

    # ==================== SLIDE 14: Closing ====================
    slide14 = prs.slides.add_slide(blank_layout)
    add_dark_background(slide14)

    # Brand
    add_text_box(slide14, 0, 2, 13.333, 1, "TALIJA",
                 font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Gold line
    line = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(3.2), Inches(1.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Quote
    add_text_box(slide14, 1.5, 3.6, 10.333, 0.6,
                 "Pozivamo vas da postanete deo naše priče.",
                 font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide14, 1.5, 4.2, 10.333, 0.5,
                 "欢迎您成为我们故事的一部分。",
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide14, 1.5, 5, 10.333, 0.6,
                 "Pravi uspeh gradi se kroz dugoročne odnose i međusobno poverenje.",
                 font_size=12, color=RgbColor(180, 180, 180), align=PP_ALIGN.CENTER)
    add_text_box(slide14, 1.5, 5.4, 10.333, 0.4,
                 "真正的成功来自长期关系与相互信任。",
                 font_size=10, color=RgbColor(140, 140, 140), align=PP_ALIGN.CENTER)

    add_text_box(slide14, 0, 6.2, 13.333, 0.4, "Hvala · 谢谢",
                 font_size=14, color=RgbColor(120, 120, 120), align=PP_ALIGN.CENTER)

    add_slide_number(slide14, 14)

    # Save presentation
    output_path = os.path.join(BASE_PATH, "TALIJA_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
